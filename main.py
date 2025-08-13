# server.py
from fastapi import FastAPI, Response
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import base64

app = FastAPI()

# Embed the three PNGs so the server is stateless
OAKLEY_LOGO_B64 = "<<<PASTE>>>" 
LINE1_B64 = "<<<PASTE>>>"
LINE2_B64 = "<<<PASTE>>>"

def write_png(path, b64s):
    with open(path, "wb") as f: f.write(base64.b64decode(b64s))

class SlideRequest(BaseModel):
    vision: str
    mission: str
    priority1: str
    priority2: str
    priority3: str
    opp1: str | None = ""
    opp2: str | None = ""
    opp3: str | None = ""

@app.post("/build", response_class=Response)
def build(req: SlideRequest):
    # write embedded assets
    write_png("Oakley logo.png", OAKLEY_LOGO_B64)
    write_png("first white line.png", LINE1_B64)
    write_png("second white line.png", LINE2_B64)

    prs = Presentation()
    prs.slide_width = Inches(26.667); prs.slide_height = Inches(15.000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    def add_text(text, l,t,w,h,font,size,rgb,bold=False,all_caps=True):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        p = tb.text_frame.paragraphs[0]; run = p.add_run()
        run.text = (text or "").upper() if all_caps else (text or "")
        f = run.font; f.name = font; f.size = Pt(size); f.color.rgb = RGBColor(*rgb); f.bold = bold

    def add_rect(l,t,w,h,fill,outline,wpt):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(*fill)
        sh.line.color.rgb = RGBColor(*outline); sh.line.width = Pt(wpt)

    # layout (same as Option A) ...
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = RGBColor(0,0,0)
    add_rect(0.85,1.90,4.82,3.02,(0,0,0),(255,255,255),1)
    add_rect(6.21,1.92,19.33,3.00,(0,0,0),(255,255,255),1)
    add_rect(0.88,5.12,24.66,2.34,(0,0,0),(255,255,255),1)
    add_rect(0.88,7.69,24.66,6.33,(0,0,0),(255,255,255),1)

    add_text("SPORTS MARKETING PORTFOLIO",0.67,0.56,10,0.6,"Avenir Next Ultra Light",16,(109,113,122))
    add_text("[INPUT] | ONE PAGE STRATEGY",0.67,0.92,15,0.8,"Avenir Next Regular",26,(255,255,255))

    add_text("VISION",1.04,2.03,4,0.5,"Helvetica",16,(255,255,255),True)
    add_text(req.vision,0.98,2.42,4.50,1.99,"Helvetica",28,(109,113,122),True)

    add_text("MISSION",6.39,2.05,4,0.5,"Helvetica",16,(255,255,255),True)
    add_text(req.mission,6.49,2.37,15,1,"Montserrat",34,(109,113,122))

    add_text("PRIORITIES",1.04,5.35,4,0.5,"Helvetica",16,(255,255,255),True)
    add_text("01",2.34,5.88,2,0.5,"Arial Black",24,(255,255,255))
    add_text("02",10.03,5.88,2,0.5,"Arial Black",24,(255,255,255))
    add_text("03",18.84,5.79,2,0.5,"Arial Black",24,(255,255,255))
    add_text(req.priority1,2.34,6.24,8,0.6,"Montserrat",24,(255,255,255),True)
    add_text(req.priority2,10.03,6.24,10,0.6,"Montserrat",24,(255,255,255),True)
    add_text(req.priority3,18.84,6.15,8,0.6,"Montserrat",24,(255,255,255),True)

    add_text("OPPORTUNITY",1.04,7.86,4,0.5,"Helvetica Neue",16,(255,255,255),True)
    add_text(req.opp1 or "",1.55,8.58,6,0.5,"Montserrat",20,(255,255,255))
    add_text(req.opp2 or "",10.21,8.40,6,0.5,"Montserrat",20,(255,255,255))
    add_text(req.opp3 or "",18.19,8.37,6,0.5,"Montserrat",20,(255,255,255))

    for (l,t) in [(1.41,10.09),(5.02,10.09),(9.78,10.03),(13.40,10.03),(18.18,10.00),(21.78,10.03)]:
        add_rect(l,t,3.27,3.17,(0,0,0),(255,255,255),1)

    slide.shapes.add_picture("Oakley logo.png", Inches(1.38), Inches(14.27), width=Inches(1.68), height=Inches(0.16))
    slide.shapes.add_picture("first white line.png", Inches(1.38), Inches(13.84), width=Inches(2.72), height=Inches(0.70))
    slide.shapes.add_picture("second white line.png", Inches(24.93), Inches(12.85), width=Inches(0.61), height=Inches(1.61))

    add_text("", 24.95, 14.00, 1, 0.5, "Arial", 12, (255,255,255))

    from io import BytesIO
    buf = BytesIO(); prs.save(buf); data = buf.getvalue()
    return Response(content=data,
                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition":"attachment; filename=one_page_strategy_with_images.pptx"})
